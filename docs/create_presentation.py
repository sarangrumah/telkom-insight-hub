"""
Generate PRESENTATION.pptx — Telkom Insight Hub
Business-level presentation with technology & modules overview.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ── Color Palette ──────────────────────────────────────────────────────────
RED_PRIMARY = RGBColor(0xE1, 0x1A, 0x27)  # Telkomsel / Komdigi red
RED_DARK = RGBColor(0xB0, 0x14, 0x1E)
DARK_BG = RGBColor(0x1A, 0x1A, 0x2E)
DARK_CARD = RGBColor(0x2D, 0x2D, 0x44)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
GOLD = RGBColor(0xFF, 0xD7, 0x00)
ACCENT_BLUE = RGBColor(0x00, 0x96, 0xD6)
ACCENT_GREEN = RGBColor(0x2E, 0xCC, 0x71)
SUBTITLE_GRAY = RGBColor(0xAA, 0xAA, 0xAA)

def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape_rect(slide, left, top, width, height, fill_color, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape

def add_rounded_rect(slide, left, top, width, height, fill_color, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    return shape

def add_text_box(slide, left, top, width, height, text, font_size=18, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Segoe UI"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_bullet_frame(slide, left, top, width, height, items, font_size=16, color=WHITE, spacing=Pt(6)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Segoe UI"
        p.space_after = spacing
        p.level = 0
    return txBox

def add_card(slide, left, top, width, height, title, items, title_color=WHITE, card_color=DARK_CARD, border=None):
    card = add_rounded_rect(slide, left, top, width, height, card_color, border)
    add_text_box(slide, left + Inches(0.3), top + Inches(0.2), width - Inches(0.6), Inches(0.5),
                 title, font_size=16, color=title_color, bold=True)
    add_bullet_frame(slide, left + Inches(0.3), top + Inches(0.65), width - Inches(0.6), height - Inches(0.85),
                     items, font_size=12, color=LIGHT_GRAY, spacing=Pt(4))
    return card

def add_accent_bar(slide):
    add_shape_rect(slide, Inches(0), Inches(7.3), prs.slide_width, Inches(0.2), RED_PRIMARY)

def add_slide_number(slide, num, total):
    add_text_box(slide, Inches(12.2), Inches(7.05), Inches(1), Inches(0.3),
                 f"{num}/{total}", font_size=10, color=SUBTITLE_GRAY, alignment=PP_ALIGN.RIGHT)

def add_section_header(slide, section_num, section_title):
    add_text_box(slide, Inches(0.8), Inches(0.3), Inches(1.5), Inches(0.5),
                 f"0{section_num}" if section_num < 10 else str(section_num),
                 font_size=48, color=RED_PRIMARY, bold=True, font_name="Segoe UI Light")
    add_text_box(slide, Inches(0.8), Inches(0.9), Inches(11), Inches(0.6),
                 section_title, font_size=32, color=WHITE, bold=True)
    add_shape_rect(slide, Inches(0.8), Inches(1.55), Inches(3), Inches(0.04), RED_PRIMARY)

TOTAL_SLIDES = 20

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title Slide
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_slide_bg(slide, DARK_BG)
add_shape_rect(slide, Inches(0), Inches(0), Inches(0.15), prs.slide_height, RED_PRIMARY)
add_text_box(slide, Inches(1.5), Inches(1.2), Inches(10), Inches(1),
             "KEMENTERIAN KOMUNIKASI DAN DIGITAL", font_size=16, color=SUBTITLE_GRAY, bold=False, font_name="Segoe UI Light")
add_text_box(slide, Inches(1.5), Inches(1.8), Inches(10), Inches(1.2),
             "Telkom Insight Hub", font_size=54, color=WHITE, bold=True)
add_text_box(slide, Inches(1.5), Inches(3.0), Inches(10), Inches(0.8),
             "Panel Manajemen Data Telekomunikasi & Analitik Strategis", font_size=24, color=GOLD, bold=False, font_name="Segoe UI Light")
add_shape_rect(slide, Inches(1.5), Inches(4.0), Inches(4), Inches(0.04), RED_PRIMARY)
add_text_box(slide, Inches(1.5), Inches(4.3), Inches(10), Inches(0.5),
             "Direktorat Jenderal Penyelenggaraan Pos dan Informatika", font_size=16, color=LIGHT_GRAY)
add_text_box(slide, Inches(1.5), Inches(4.8), Inches(10), Inches(0.5),
             "Presentasi Pengembangan Sistem — April 2026", font_size=14, color=SUBTITLE_GRAY)
# Version badge
badge = add_rounded_rect(slide, Inches(1.5), Inches(5.5), Inches(2), Inches(0.45), RED_PRIMARY)
add_text_box(slide, Inches(1.5), Inches(5.53), Inches(2), Inches(0.4),
             "Version 1.0 — Production", font_size=12, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_accent_bar(slide)
add_slide_number(slide, 1, TOTAL_SLIDES)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Agenda
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_section_header(slide, 0, "AGENDA")

agenda_items = [
    ("01", "Latar Belakang & Tujuan Strategis"),
    ("02", "Arsitektur Sistem & Teknologi"),
    ("03", "Modul Inti — Manajemen Data Telekomunikasi"),
    ("04", "Modul Analitik — BPS & Potensi Pasar"),
    ("05", "Modul Pendukung — Registrasi, Tiket, FAQ"),
    ("06", "Integrasi Sistem Eksternal"),
    ("07", "Keamanan & Kepatuhan (DevSecOps)"),
    ("08", "Dashboard & Visualisasi Data"),
    ("09", "Capaian & Progress Pengembangan"),
    ("10", "Roadmap & Rencana Selanjutnya"),
]
for i, (num, title) in enumerate(agenda_items):
    y = Inches(2.0) + Inches(0.5) * i
    add_rounded_rect(slide, Inches(1.0), y, Inches(0.7), Inches(0.4), RED_PRIMARY if i < 5 else DARK_CARD, RED_PRIMARY if i >= 5 else None)
    add_text_box(slide, Inches(1.0), y + Inches(0.02), Inches(0.7), Inches(0.35),
                 num, font_size=14, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1.9), y + Inches(0.02), Inches(8), Inches(0.35),
                 title, font_size=16, color=WHITE)

add_accent_bar(slide)
add_slide_number(slide, 2, TOTAL_SLIDES)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Latar Belakang
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_section_header(slide, 1, "LATAR BELAKANG & TUJUAN STRATEGIS")

add_text_box(slide, Inches(0.8), Inches(1.9), Inches(11), Inches(0.5),
             "Tantangan Utama", font_size=20, color=GOLD, bold=True)
challenges = [
    "Data perizinan telekomunikasi tersebar di berbagai sistem (POSTEL, OSS, SDPPI)",
    "Belum ada platform terpadu untuk analisis potensi pasar telekomunikasi",
    "Proses verifikasi perusahaan masih manual dan memakan waktu",
    "Minimnya integrasi data statistik (BPS) dengan data perizinan",
    "Kebutuhan dashboard real-time untuk pengambilan keputusan strategis",
]
add_bullet_frame(slide, Inches(1.0), Inches(2.4), Inches(5.5), Inches(3),
                 [f"\u2022  {c}" for c in challenges], font_size=13, color=LIGHT_GRAY)

add_text_box(slide, Inches(7.0), Inches(1.9), Inches(6), Inches(0.5),
             "Tujuan Pengembangan", font_size=20, color=GOLD, bold=True)
goals = [
    "Menyediakan Single Source of Truth data telekomunikasi",
    "Mengintegrasikan data BPS untuk analisis potensi pasar",
    "Otomasi proses verifikasi & registrasi perusahaan",
    "Dashboard analitik untuk mendukung kebijakan publik",
    "Integrasi seamless dengan sistem e-Telekomunikasi",
]
add_bullet_frame(slide, Inches(7.2), Inches(2.4), Inches(5.5), Inches(3),
                 [f"\u2022  {g}" for g in goals], font_size=13, color=LIGHT_GRAY)

# Key metrics boxes
metrics = [("8+", "Jenis\nLayanan"), ("34+", "Provinsi\nCakupan"), ("500+", "Tabel\nDatabase"), ("24/7", "Monitoring\nReal-time")]
for i, (val, label) in enumerate(metrics):
    x = Inches(0.8) + Inches(3.05) * i
    card = add_rounded_rect(slide, x, Inches(5.5), Inches(2.8), Inches(1.3), DARK_CARD, RED_PRIMARY)
    add_text_box(slide, x, Inches(5.6), Inches(2.8), Inches(0.6),
                 val, font_size=36, color=RED_PRIMARY, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x, Inches(6.15), Inches(2.8), Inches(0.6),
                 label, font_size=12, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

add_accent_bar(slide)
add_slide_number(slide, 3, TOTAL_SLIDES)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 4 — System Architecture Overview
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_section_header(slide, 2, "ARSITEKTUR SISTEM")

# 3-tier architecture boxes
tiers = [
    ("FRONTEND (Client Layer)", ACCENT_BLUE, [
        "React 18 + TypeScript 5.8",
        "TailwindCSS + shadcn/ui Design System",
        "Mapbox GL untuk Visualisasi Peta",
        "Recharts untuk Grafik & Analitik",
        "React Query untuk State Management",
    ]),
    ("BACKEND (API Layer)", RED_PRIMARY, [
        "Node.js + Express 4.19",
        "JWT Authentication + RBAC",
        "WebSocket Real-time Updates",
        "RESTful API — 60+ Endpoints",
        "Background Job Scheduler (node-cron)",
    ]),
    ("DATABASE & INTEGRATION", ACCENT_GREEN, [
        "PostgreSQL — Relational Database",
        "40+ Tables dengan Indexing Optimal",
        "BPS API Integration (webapi.bps.go.id)",
        "KOMINFO Tariff API Sync",
        "e-Telekomunikasi Service Bridge",
    ]),
]

for i, (title, color, items) in enumerate(tiers):
    x = Inches(0.5) + Inches(4.1) * i
    header = add_rounded_rect(slide, x, Inches(1.9), Inches(3.9), Inches(0.55), color)
    add_text_box(slide, x, Inches(1.95), Inches(3.9), Inches(0.45),
                 title, font_size=14, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    card = add_rounded_rect(slide, x, Inches(2.5), Inches(3.9), Inches(3.2), DARK_CARD)
    add_bullet_frame(slide, x + Inches(0.2), Inches(2.65), Inches(3.5), Inches(3),
                     [f"\u25B8  {item}" for item in items], font_size=12, color=LIGHT_GRAY, spacing=Pt(8))

# Deployment info
add_text_box(slide, Inches(0.8), Inches(6.0), Inches(11), Inches(0.4),
             "Deployment: Docker Container  |  PM2 Process Manager  |  Apache/Nginx Reverse Proxy  |  Systemd Service",
             font_size=12, color=SUBTITLE_GRAY, alignment=PP_ALIGN.CENTER)

add_accent_bar(slide)
add_slide_number(slide, 4, TOTAL_SLIDES)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Technology Stack Detail
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_section_header(slide, 2, "TECHNOLOGY STACK")

categories = [
    ("Frontend Framework", ["React 18.3.1", "TypeScript 5.8", "Vite 5.4 Build Tool", "React Router v6"]),
    ("UI & Design", ["TailwindCSS 3.4", "shadcn/ui (47 components)", "Lucide Icons", "Recharts 2.15"]),
    ("State & Data", ["TanStack React Query v5", "React Hook Form + Zod", "WebSocket Client", "Axios HTTP Client"]),
    ("Backend", ["Express 4.19", "PostgreSQL (pg 8.13)", "JWT + bcrypt Auth", "Multer File Upload"]),
    ("Security", ["Helmet CSP Headers", "CORS Protection", "Rate Limiting", "ESLint Security Plugin"]),
    ("DevOps", ["Docker + Systemd", "Vitest Unit Testing", "Husky Pre-commit", "PM2 Process Mgr"]),
]

for i, (cat_title, items) in enumerate(categories):
    col = i % 3
    row = i // 3
    x = Inches(0.5) + Inches(4.1) * col
    y = Inches(1.9) + Inches(2.6) * row
    add_rounded_rect(slide, x, y, Inches(3.9), Inches(2.4), DARK_CARD)
    add_text_box(slide, x + Inches(0.2), y + Inches(0.1), Inches(3.5), Inches(0.4),
                 cat_title, font_size=15, color=ACCENT_BLUE, bold=True)
    add_bullet_frame(slide, x + Inches(0.3), y + Inches(0.55), Inches(3.3), Inches(1.7),
                     [f"\u25CF  {it}" for it in items], font_size=12, color=LIGHT_GRAY, spacing=Pt(5))

add_accent_bar(slide)
add_slide_number(slide, 5, TOTAL_SLIDES)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Module Overview Map
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_section_header(slide, 3, "PETA MODUL SISTEM")

modules = [
    ("Manajemen Data\nTelekomunikasi", RED_PRIMARY, "MODUL INTI"),
    ("Registrasi &\nVerifikasi Perusahaan", RED_PRIMARY, "MODUL INTI"),
    ("BPS Statistical\nData Integration", ACCENT_BLUE, "ANALITIK"),
    ("Telecom Market\nPotential Scoring", ACCENT_BLUE, "ANALITIK"),
    ("Dashboard &\nVisualisasi", ACCENT_BLUE, "ANALITIK"),
    ("Support &\nTicketing", ACCENT_GREEN, "PENDUKUNG"),
    ("FAQ\nManagement", ACCENT_GREEN, "PENDUKUNG"),
    ("User & Permission\nManagement (RBAC)", ACCENT_GREEN, "PENDUKUNG"),
    ("Integrasi Sistem\nEksternal", GOLD, "INTEGRASI"),
    ("DevSecOps &\nMonitoring", GOLD, "INTEGRASI"),
]

for i, (name, color, category) in enumerate(modules):
    col = i % 5
    row = i // 5
    x = Inches(0.3) + Inches(2.55) * col
    y = Inches(2.0) + Inches(2.7) * row
    card = add_rounded_rect(slide, x, y, Inches(2.35), Inches(2.2), DARK_CARD, color)
    # Category badge
    badge_w = Inches(1.6)
    add_rounded_rect(slide, x + Inches(0.37), y + Inches(0.15), badge_w, Inches(0.35), color)
    add_text_box(slide, x + Inches(0.37), y + Inches(0.17), badge_w, Inches(0.3),
                 category, font_size=9, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.15), y + Inches(0.7), Inches(2.05), Inches(1.2),
                 name, font_size=14, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

add_accent_bar(slide)
add_slide_number(slide, 6, TOTAL_SLIDES)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Manajemen Data Telekomunikasi
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_section_header(slide, 3, "MODUL: MANAJEMEN DATA TELEKOMUNIKASI")

add_text_box(slide, Inches(0.8), Inches(1.8), Inches(11), Inches(0.5),
             "Pengelolaan terpadu seluruh data perizinan telekomunikasi di Indonesia",
             font_size=16, color=SUBTITLE_GRAY)

services = [
    ("Jasa Telekomunikasi", "Penyelenggaraan jasa\ntelekomunikasi umum"),
    ("Jaringan Telekomunikasi", "Penyelenggaraan jaringan\ntelekomunikasi"),
    ("Telekomunikasi Khusus", "Penyelenggaraan telekomunikasi\nkhusus (Telsus)"),
    ("ISR", "Izin Stasiun Radio\nuntuk spektrum frekuensi"),
    ("SKLO", "Sertifikat Kelaikan\nLayak Operasi"),
    ("LKO", "Laporan Kelaikan\nOperasi perangkat"),
    ("Penomoran", "Pengelolaan alokasi\nnomor telekomunikasi"),
    ("Tarif", "Pengaturan tarif\nlayanan telekomunikasi"),
]

for i, (svc, desc) in enumerate(services):
    col = i % 4
    row = i // 4
    x = Inches(0.5) + Inches(3.1) * col
    y = Inches(2.4) + Inches(2.0) * row
    add_rounded_rect(slide, x, y, Inches(2.9), Inches(1.7), DARK_CARD, RED_PRIMARY)
    add_text_box(slide, x + Inches(0.15), y + Inches(0.15), Inches(2.6), Inches(0.4),
                 svc, font_size=14, color=RED_PRIMARY, bold=True)
    add_text_box(slide, x + Inches(0.15), y + Inches(0.6), Inches(2.6), Inches(0.9),
                 desc, font_size=11, color=LIGHT_GRAY)

# Feature highlights
add_text_box(slide, Inches(0.8), Inches(6.5), Inches(11), Inches(0.4),
             "Fitur: CRUD Operations  |  Excel Import/Export  |  Advanced Search & Filter  |  Location-based Filtering  |  Pagination & Sorting",
             font_size=12, color=ACCENT_BLUE, alignment=PP_ALIGN.CENTER)

add_accent_bar(slide)
add_slide_number(slide, 7, TOTAL_SLIDES)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Company Registration & Verification
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_section_header(slide, 3, "MODUL: REGISTRASI & VERIFIKASI PERUSAHAAN")

# Workflow steps
steps = [
    ("1", "REGISTRASI", "Pelaku usaha mengisi\ndata perusahaan &\nmengunggah dokumen\n(NIB, NPWP, Akta, KTP)"),
    ("2", "SUBMITTED", "Sistem menerima\npendaftaran dan\nmembuat akun\ndengan status pending"),
    ("3", "IN REVIEW", "Admin memverifikasi\nkelengkapan dokumen\ndan keabsahan data\nperusahaan"),
    ("4", "VERIFIED", "Perusahaan terverifikasi\ndan mendapat akses\npenuh ke seluruh\nfitur sistem"),
]

for i, (num, title, desc) in enumerate(steps):
    x = Inches(0.5) + Inches(3.2) * i
    # Step number circle
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(1.1), Inches(2.0), Inches(0.7), Inches(0.7))
    circle.fill.solid()
    circle.fill.fore_color.rgb = RED_PRIMARY
    circle.line.fill.background()
    add_text_box(slide, x + Inches(1.1), Inches(2.05), Inches(0.7), Inches(0.6),
                 num, font_size=24, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x, Inches(2.85), Inches(2.9), Inches(0.4),
                 title, font_size=16, color=GOLD, bold=True, alignment=PP_ALIGN.CENTER)
    add_rounded_rect(slide, x, Inches(3.3), Inches(2.9), Inches(2.0), DARK_CARD)
    add_text_box(slide, x + Inches(0.15), Inches(3.45), Inches(2.6), Inches(1.7),
                 desc, font_size=12, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)
    # Arrow between steps
    if i < 3:
        arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x + Inches(3.0), Inches(2.15), Inches(0.3), Inches(0.4))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = RED_PRIMARY
        arrow.line.fill.background()

# Additional features
add_card(slide, Inches(0.5), Inches(5.6), Inches(3.8), Inches(1.4),
         "Dokumen yang Dikelola", [
             "\u25B8 NIB (Nomor Induk Berusaha)",
             "\u25B8 NPWP Perusahaan",
             "\u25B8 Akta Pendirian",
             "\u25B8 KTP Penanggung Jawab",
         ], title_color=ACCENT_BLUE)

add_card(slide, Inches(4.6), Inches(5.6), Inches(3.8), Inches(1.4),
         "Fitur Verifikasi Admin", [
             "\u25B8 Dashboard verifikasi terpusat",
             "\u25B8 Catatan koreksi & feedback",
             "\u25B8 Status tracking real-time",
             "\u25B8 Bulk action management",
         ], title_color=ACCENT_BLUE)

add_card(slide, Inches(8.7), Inches(5.6), Inches(4.1), Inches(1.4),
         "Status Perusahaan", [
             "\u25B8 Submitted \u2192 In Review \u2192 Verified",
             "\u25B8 Needs Correction (with notes)",
             "\u25B8 Rejected (with reason)",
             "\u25B8 PIC Management per company",
         ], title_color=ACCENT_BLUE)

add_accent_bar(slide)
add_slide_number(slide, 8, TOTAL_SLIDES)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 9 — BPS Integration
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_section_header(slide, 4, "MODUL: INTEGRASI DATA BPS")

add_text_box(slide, Inches(0.8), Inches(1.8), Inches(11), Inches(0.5),
             "Integrasi langsung dengan Badan Pusat Statistik (BPS) untuk analisis data makroekonomi & telekomunikasi",
             font_size=14, color=SUBTITLE_GRAY)

# BPS Features
bps_cards = [
    ("Konfigurasi API", ACCENT_BLUE, [
        "API Key Management",
        "Rate Limit Configuration",
        "Base URL Setting",
        "Health Check Monitoring",
    ]),
    ("Manajemen Area", ACCENT_GREEN, [
        "Seleksi Provinsi & Kabupaten",
        "34 Provinsi Indonesia",
        "Monitoring area terpilih",
        "Filter data per wilayah",
    ]),
    ("Variabel Statistik", GOLD, [
        "Search katalog BPS",
        "Registrasi variabel aktif",
        "Kategori & filter variabel",
        "Toggle status aktif/nonaktif",
    ]),
    ("Sinkronisasi Data", RED_PRIMARY, [
        "Full & Incremental Sync",
        "Sync History Tracking",
        "Background Job Scheduler",
        "Error Handling & Retry",
    ]),
]

for i, (title, color, items) in enumerate(bps_cards):
    x = Inches(0.3) + Inches(3.2) * i
    add_rounded_rect(slide, x, Inches(2.4), Inches(3.0), Inches(0.5), color)
    add_text_box(slide, x, Inches(2.43), Inches(3.0), Inches(0.45),
                 title, font_size=14, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_rounded_rect(slide, x, Inches(3.0), Inches(3.0), Inches(2.2), DARK_CARD, color)
    add_bullet_frame(slide, x + Inches(0.2), Inches(3.15), Inches(2.6), Inches(2.0),
                     [f"\u25B8  {it}" for it in items], font_size=12, color=LIGHT_GRAY, spacing=Pt(8))

# Data flow
add_text_box(slide, Inches(0.8), Inches(5.5), Inches(11.5), Inches(0.5),
             "Data Pipeline", font_size=18, color=GOLD, bold=True)

pipeline = ["webapi.bps.go.id", "API Fetcher Service", "Data Normalizer", "PostgreSQL Storage", "Visualization Layer"]
for i, step in enumerate(pipeline):
    x = Inches(0.3) + Inches(2.55) * i
    color = [ACCENT_BLUE, RED_PRIMARY, ACCENT_GREEN, GOLD, ACCENT_BLUE][i]
    add_rounded_rect(slide, x, Inches(6.0), Inches(2.3), Inches(0.6), color)
    add_text_box(slide, x, Inches(6.05), Inches(2.3), Inches(0.5),
                 step, font_size=11, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    if i < 4:
        arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x + Inches(2.35), Inches(6.1), Inches(0.2), Inches(0.35))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = LIGHT_GRAY
        arrow.line.fill.background()

add_accent_bar(slide)
add_slide_number(slide, 9, TOTAL_SLIDES)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Telecom Market Potential
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_section_header(slide, 4, "MODUL: POTENSI PASAR TELEKOMUNIKASI")

add_text_box(slide, Inches(0.8), Inches(1.8), Inches(11), Inches(0.5),
             "Algoritma scoring multi-dimensi untuk mengukur potensi pasar telekomunikasi per wilayah",
             font_size=14, color=SUBTITLE_GRAY)

# Scoring components
add_card(slide, Inches(0.5), Inches(2.3), Inches(3.8), Inches(2.8),
         "Market Activity Score", [
             "\u25B8 License Count per Area",
             "\u25B8 Service Diversity Index",
             "\u25B8 Active vs Total License Ratio",
             "\u25B8 Historical Growth Trend",
             "\u25B8 Weighted by service type:",
             "   Jasa 20%, Jaringan 25%,",
             "   ISR 15%, Telsus 15%, etc.",
         ], title_color=RED_PRIMARY, border=RED_PRIMARY)

add_card(slide, Inches(4.7), Inches(2.3), Inches(3.8), Inches(2.8),
         "Untapped Opportunity Score", [
             "\u25B8 Population Density Analysis",
             "\u25B8 Infrastructure Gap Assessment",
             "\u25B8 Under-served Area Detection",
             "\u25B8 Growth Potential Calculation",
             "\u25B8 Demographic Indicators",
             "\u25B8 Economic Activity Proxies",
         ], title_color=ACCENT_BLUE, border=ACCENT_BLUE)

add_card(slide, Inches(8.9), Inches(2.3), Inches(3.9), Inches(2.8),
         "BPS Demand Score (V2)", [
             "\u25B8 BPS Statistical Blending",
             "\u25B8 GDP per Capita Integration",
             "\u25B8 Internet Penetration Rate",
             "\u25B8 Household ICT Expenditure",
             "\u25B8 Education & Urbanization",
             "\u25B8 Configurable Weight System",
         ], title_color=ACCENT_GREEN, border=ACCENT_GREEN)

# Output
add_text_box(slide, Inches(0.8), Inches(5.4), Inches(11), Inches(0.4),
             "Output & Visualisasi", font_size=18, color=GOLD, bold=True)

outputs = [
    ("Tiered Ranking\nA / B / C / D", RED_PRIMARY),
    ("Interactive Map\n(Mapbox GL)", ACCENT_BLUE),
    ("Radar Chart\nMulti-dimensi", ACCENT_GREEN),
    ("Sortable Table\nwith Scores", GOLD),
]
for i, (label, color) in enumerate(outputs):
    x = Inches(0.5) + Inches(3.2) * i
    add_rounded_rect(slide, x, Inches(5.9), Inches(2.9), Inches(1.0), DARK_CARD, color)
    add_text_box(slide, x, Inches(6.0), Inches(2.9), Inches(0.8),
                 label, font_size=13, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

add_accent_bar(slide)
add_slide_number(slide, 10, TOTAL_SLIDES)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 11 — Dashboard & Visualization
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_section_header(slide, 4, "DASHBOARD & VISUALISASI DATA")

dashboards = [
    ("Main Dashboard", [
        "KPI Cards (total data per service)",
        "Time-series trend analysis",
        "Service distribution breakdown",
        "Recent activity feed",
    ]),
    ("BPS Data Visualization", [
        "Line charts — time-series trends",
        "Bar charts — area comparison",
        "Multi-variable overlay",
        "Sync history & status",
    ]),
    ("Telecom Potential Map", [
        "Color-coded Mapbox choropleth",
        "Click-to-detail interaction",
        "Radar chart per area",
        "Pie chart service breakdown",
    ]),
    ("Admin Analytics", [
        "Company verification funnel",
        "User activity monitoring",
        "Ticket SLA performance",
        "API health & latency",
    ]),
]

for i, (title, items) in enumerate(dashboards):
    col = i % 2
    row = i // 2
    x = Inches(0.5) + Inches(6.3) * col
    y = Inches(1.9) + Inches(2.6) * row
    add_card(slide, x, y, Inches(5.9), Inches(2.3), title, [f"\u25B8  {it}" for it in items],
             title_color=ACCENT_BLUE, border=ACCENT_BLUE)

add_accent_bar(slide)
add_slide_number(slide, 11, TOTAL_SLIDES)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 12 — Support & FAQ Modules
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_section_header(slide, 5, "MODUL PENDUKUNG: SUPPORT & FAQ")

# Ticketing
add_card(slide, Inches(0.5), Inches(1.9), Inches(5.9), Inches(3.0),
         "Support & Ticketing System", [
             "\u25B8  Pembuatan tiket dari antarmuka publik",
             "\u25B8  Assignment tiket ke admin yang tersedia",
             "\u25B8  Real-time messaging via WebSocket",
             "\u25B8  Status tracking: Open \u2192 In Progress \u2192 Resolved \u2192 Closed",
             "\u25B8  Admin dashboard dengan filter & search",
             "\u25B8  Unread message counter real-time",
             "\u25B8  SLA metrics tracking & reporting",
             "\u25B8  Conversation threading per tiket",
         ], title_color=RED_PRIMARY, border=RED_PRIMARY)

# FAQ
add_card(slide, Inches(6.8), Inches(1.9), Inches(5.9), Inches(3.0),
         "FAQ Management", [
             "\u25B8  Public FAQ viewer dengan kategori",
             "\u25B8  Admin FAQ editor (CRUD)",
             "\u25B8  Kategori management",
             "\u25B8  Search functionality",
             "\u25B8  API-backed content management",
             "\u25B8  Responsive mobile-friendly layout",
         ], title_color=ACCENT_GREEN, border=ACCENT_GREEN)

# RBAC
add_card(slide, Inches(0.5), Inches(5.2), Inches(12.2), Inches(1.8),
         "User & Permission Management (RBAC)", [
             "\u25B8  6 Role: super_admin, internal_admin, pelaku_usaha, pengolah_data, internal_group, guest",
             "\u25B8  Module-level permissions (read/create/update/delete) dan field-level granularity",
             "\u25B8  Admin interface untuk assignment permission    |    PermissionGuard component untuk UI protection",
         ], title_color=GOLD, border=GOLD)

add_accent_bar(slide)
add_slide_number(slide, 12, TOTAL_SLIDES)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 13 — External Integrations
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_section_header(slide, 6, "INTEGRASI SISTEM EKSTERNAL")

integrations = [
    ("e-Telekomunikasi", "Platform perizinan utama Komdigi", RED_PRIMARY, [
        "Service-to-service API key auth",
        "Sinkronisasi data perizinan",
        "Shared authentication bridge",
        "Real-time data exchange",
    ]),
    ("BPS (webapi.bps.go.id)", "Badan Pusat Statistik", ACCENT_BLUE, [
        "Statistical data fetching",
        "Variable catalog search",
        "Time-series data sync",
        "Rate-limited API calls",
    ]),
    ("KOMINFO Tarif", "tariftel.komdigi.go.id", ACCENT_GREEN, [
        "Tariff data synchronization",
        "Periodic auto-sync (cron)",
        "Manual sync trigger",
        "History & error tracking",
    ]),
    ("OSS / POSTEL / SDPPI", "Sistem Perizinan Nasional", GOLD, [
        "OSS: Business registration",
        "POSTEL: Legacy licensing",
        "SDPPI: Radio spectrum data",
        "Pluggable adapter pattern",
    ]),
]

for i, (name, subtitle, color, items) in enumerate(integrations):
    col = i % 2
    row = i // 2
    x = Inches(0.5) + Inches(6.3) * col
    y = Inches(1.9) + Inches(2.6) * row
    add_rounded_rect(slide, x, y, Inches(5.9), Inches(2.3), DARK_CARD, color)
    add_text_box(slide, x + Inches(0.3), y + Inches(0.15), Inches(3.5), Inches(0.4),
                 name, font_size=16, color=color, bold=True)
    add_text_box(slide, x + Inches(0.3), y + Inches(0.5), Inches(3.5), Inches(0.3),
                 subtitle, font_size=11, color=SUBTITLE_GRAY)
    add_bullet_frame(slide, x + Inches(0.3), y + Inches(0.85), Inches(5.3), Inches(1.3),
                     [f"\u25B8  {it}" for it in items], font_size=12, color=LIGHT_GRAY, spacing=Pt(5))

add_accent_bar(slide)
add_slide_number(slide, 13, TOTAL_SLIDES)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 14 — Security & DevSecOps
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_section_header(slide, 7, "KEAMANAN & KEPATUHAN (DevSecOps)")

sec_areas = [
    ("Authentication & Authorization", RED_PRIMARY, [
        "JWT Token + Refresh Rotation",
        "bcrypt Password Hashing",
        "RBAC 6-level Role System",
        "Session Management (httpOnly cookie)",
        "External Auth Integration",
    ]),
    ("Network & Transport", ACCENT_BLUE, [
        "Helmet Security Headers (CSP, HSTS)",
        "CORS Whitelist Configuration",
        "Rate Limiting (express-rate-limit)",
        "HTTPS/TLS Enforcement",
        "WebSocket Secure (WSS)",
    ]),
    ("Data Protection", ACCENT_GREEN, [
        "Parameterized SQL Queries",
        "Input Sanitization (XSS)",
        "File Upload Type Validation",
        "Sensitive Data Encryption",
        "Audit Trail Logging",
    ]),
    ("Monitoring & Compliance", GOLD, [
        "Security Metrics Dashboard",
        "API Call Audit Logging",
        "Login Attempt Tracking",
        "Performance Monitoring",
        "ESLint Security Plugin",
    ]),
]

for i, (title, color, items) in enumerate(sec_areas):
    col = i % 2
    row = i // 2
    x = Inches(0.5) + Inches(6.3) * col
    y = Inches(1.9) + Inches(2.4) * row
    add_card(slide, x, y, Inches(5.9), Inches(2.2), title, [f"\u25B8  {it}" for it in items],
             title_color=color, border=color)

add_accent_bar(slide)
add_slide_number(slide, 14, TOTAL_SLIDES)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 15 — User Roles & Access Control
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_section_header(slide, 7, "ROLE-BASED ACCESS CONTROL (RBAC)")

roles = [
    ("Super Admin", "Akses penuh ke seluruh sistem, konfigurasi, dan manajemen pengguna", RED_PRIMARY),
    ("Internal Admin", "Verifikasi perusahaan, manajemen tiket, pengelolaan data", RED_DARK),
    ("Pelaku Usaha", "Registrasi perusahaan, melihat data perizinan, support", ACCENT_BLUE),
    ("Pengolah Data", "Import/export data, CRUD telekom data, visualisasi", ACCENT_GREEN),
    ("Internal Group", "Akses terbatas ke data internal dan dashboard", GOLD),
    ("Guest", "Akses publik: search data, FAQ, dan registrasi", SUBTITLE_GRAY),
]

for i, (role, desc, color) in enumerate(roles):
    y = Inches(1.9) + Inches(0.85) * i
    add_rounded_rect(slide, Inches(0.5), y, Inches(2.5), Inches(0.7), color)
    add_text_box(slide, Inches(0.5), y + Inches(0.1), Inches(2.5), Inches(0.5),
                 role, font_size=14, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(3.2), y + Inches(0.1), Inches(9.5), Inches(0.5),
                 desc, font_size=14, color=LIGHT_GRAY)

add_accent_bar(slide)
add_slide_number(slide, 15, TOTAL_SLIDES)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 16 — API Endpoints Overview
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_section_header(slide, 8, "API ENDPOINTS — 60+ REST ENDPOINTS")

api_groups = [
    ("Authentication", "6 endpoints", ["/auth/register", "/auth/login", "/auth/refresh", "/auth/logout", "/auth/check-email", "/auth/login-etelekomunikasi"]),
    ("Telekom Data", "5 endpoints", ["/telekom-data (CRUD)", "/telekom-data/:id", "/telekom-data/excel", "Search & Filter", "Pagination"]),
    ("BPS Integration", "12 endpoints", ["/bps/config", "/bps/areas", "/bps/variables", "/bps/data", "/bps/sync/*", "/bps/catalog/search"]),
    ("Telecom Potential", "10 endpoints", ["/telecom-potential/scores", "/telecom-potential/v2/*", "/telecom-potential/sync", "GeoJSON export", "Area detail"]),
    ("Admin Operations", "8 endpoints", ["/admin/companies", "/admin/users", "/admin/permissions", "/admin/verification", "Bulk operations"]),
    ("Support & FAQ", "8 endpoints", ["/tickets (CRUD)", "/tickets/:id/messages", "/tickets/:id/assign", "/faqs (CRUD)", "Real-time WS"]),
    ("Integrations", "5 endpoints", ["/integrations", "/integrations/:name", "/integrations/sync", "/kominfo-sync/*", "Status & History"]),
    ("Services", "8 endpoints", ["/jasa", "/jaringan", "/telsus", "/isr", "/sklo", "/lko", "/penomoran", "/tarif"]),
]

for i, (group, count, endpoints) in enumerate(api_groups):
    col = i % 4
    row = i // 4
    x = Inches(0.3) + Inches(3.2) * col
    y = Inches(1.9) + Inches(2.7) * row
    add_rounded_rect(slide, x, y, Inches(3.0), Inches(2.5), DARK_CARD, ACCENT_BLUE)
    add_text_box(slide, x + Inches(0.15), y + Inches(0.1), Inches(2.0), Inches(0.35),
                 group, font_size=13, color=ACCENT_BLUE, bold=True)
    badge = add_rounded_rect(slide, x + Inches(2.0), y + Inches(0.12), Inches(0.85), Inches(0.28), ACCENT_BLUE)
    add_text_box(slide, x + Inches(2.0), y + Inches(0.12), Inches(0.85), Inches(0.28),
                 count, font_size=9, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_bullet_frame(slide, x + Inches(0.15), y + Inches(0.5), Inches(2.7), Inches(1.9),
                     endpoints, font_size=10, color=LIGHT_GRAY, spacing=Pt(3))

add_accent_bar(slide)
add_slide_number(slide, 16, TOTAL_SLIDES)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 17 — Database Architecture
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_section_header(slide, 8, "DATABASE ARCHITECTURE — 40+ TABLES")

db_groups = [
    ("User & Company", RED_PRIMARY, ["profiles", "user_roles", "companies", "company_documents", "person_in_charge", "pic_documents"]),
    ("Telekom Data", ACCENT_BLUE, ["telekom_data", "services", "sub_services", "license_services", "provinces", "kabupaten"]),
    ("BPS Statistical", ACCENT_GREEN, ["bps_config", "bps_variables", "bps_monitored_areas", "bps_statistical_data", "bps_sync_history", "bps_api_requests"]),
    ("Support & Audit", GOLD, ["tickets", "ticket_messages", "ticket_sla_metrics", "audit_logs", "activity_logs", "api_integration_logs"]),
    ("Permissions", RGBColor(0xE9, 0x1E, 0x63), ["permissions", "permission_templates", "record_permissions", "modules", "fields", "login_attempts"]),
    ("Sync & Config", RGBColor(0x9C, 0x27, 0xB0), ["sync_configurations", "kominfo_tarif_data", "sync_status", "telecom_potential_config", "telecom_potential_area_scores", "faq_categories"]),
]

for i, (group, color, tables) in enumerate(db_groups):
    col = i % 3
    row = i // 3
    x = Inches(0.3) + Inches(4.2) * col
    y = Inches(1.9) + Inches(2.7) * row
    add_rounded_rect(slide, x, y, Inches(4.0), Inches(2.5), DARK_CARD, color)
    add_text_box(slide, x + Inches(0.2), y + Inches(0.1), Inches(3.6), Inches(0.35),
                 group, font_size=14, color=color, bold=True)
    add_bullet_frame(slide, x + Inches(0.2), y + Inches(0.5), Inches(3.6), Inches(1.8),
                     [f"\u25B8  {t}" for t in tables], font_size=11, color=LIGHT_GRAY, spacing=Pt(4))

add_accent_bar(slide)
add_slide_number(slide, 17, TOTAL_SLIDES)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 18 — Development Progress
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_section_header(slide, 9, "CAPAIAN & PROGRESS PENGEMBANGAN")

# Milestones
milestones = [
    ("Q1 2026", "Foundation", "100%", [
        "Core architecture setup",
        "Authentication & RBAC",
        "Telekom data CRUD",
        "Database schema design",
    ]),
    ("Q2 2026", "Integration", "85%", [
        "BPS API integration",
        "KOMINFO tariff sync",
        "Company verification",
        "Excel import/export",
    ]),
    ("Q3 2026", "Analytics", "70%", [
        "Market potential scoring",
        "Interactive map dashboard",
        "Data visualization charts",
        "Real-time WebSocket",
    ]),
    ("Q4 2026", "Optimization", "Planned", [
        "Performance tuning",
        "Advanced analytics",
        "Mobile responsive",
        "UAT & production deploy",
    ]),
]

for i, (period, phase, progress, items) in enumerate(milestones):
    x = Inches(0.3) + Inches(3.2) * i
    # Progress indicator
    color = ACCENT_GREEN if progress == "100%" else ACCENT_BLUE if "%" in progress else GOLD
    add_rounded_rect(slide, x, Inches(1.9), Inches(3.0), Inches(0.5), color)
    add_text_box(slide, x, Inches(1.93), Inches(3.0), Inches(0.45),
                 f"{period}  —  {progress}", font_size=14, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x, Inches(2.5), Inches(3.0), Inches(0.4),
                 phase, font_size=16, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    add_rounded_rect(slide, x, Inches(2.95), Inches(3.0), Inches(2.2), DARK_CARD, color)
    add_bullet_frame(slide, x + Inches(0.2), Inches(3.1), Inches(2.6), Inches(2.0),
                     [f"\u2713  {it}" if progress == "100%" else f"\u25B8  {it}" for it in items],
                     font_size=12, color=LIGHT_GRAY, spacing=Pt(8))

# Stats summary
stats = [
    ("22+", "Pages/Views"),
    ("55+", "Components"),
    ("60+", "API Endpoints"),
    ("40+", "DB Tables"),
    ("13", "Router Files"),
    ("8", "Service Types"),
]
for i, (val, label) in enumerate(stats):
    x = Inches(0.3) + Inches(2.15) * i
    add_rounded_rect(slide, x, Inches(5.5), Inches(2.0), Inches(1.2), DARK_CARD, RED_PRIMARY)
    add_text_box(slide, x, Inches(5.6), Inches(2.0), Inches(0.55),
                 val, font_size=28, color=RED_PRIMARY, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x, Inches(6.1), Inches(2.0), Inches(0.4),
                 label, font_size=12, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

add_accent_bar(slide)
add_slide_number(slide, 18, TOTAL_SLIDES)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 19 — Roadmap
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_section_header(slide, 10, "ROADMAP & RENCANA SELANJUTNYA")

roadmap = [
    ("Short Term (Q2 2026)", ACCENT_GREEN, [
        "Finalisasi integrasi BPS & KOMINFO",
        "Completion company verification workflow",
        "UAT testing 528 skenario",
        "Performance optimization",
        "Documentation completion",
    ]),
    ("Mid Term (Q3 2026)", ACCENT_BLUE, [
        "Advanced analytics & reporting",
        "Mobile responsive optimization",
        "API rate limiting enhancement",
        "Batch processing for large datasets",
        "Integration with OSS/POSTEL/SDPPI",
    ]),
    ("Long Term (Q4 2026+)", GOLD, [
        "AI-powered anomaly detection",
        "Predictive analytics for market trends",
        "Multi-tenant architecture",
        "National deployment & scalability",
        "Open data portal for public access",
    ]),
]

for i, (period, color, items) in enumerate(roadmap):
    x = Inches(0.3) + Inches(4.2) * i
    add_rounded_rect(slide, x, Inches(1.9), Inches(4.0), Inches(0.55), color)
    add_text_box(slide, x, Inches(1.93), Inches(4.0), Inches(0.5),
                 period, font_size=16, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_rounded_rect(slide, x, Inches(2.55), Inches(4.0), Inches(3.5), DARK_CARD, color)
    add_bullet_frame(slide, x + Inches(0.2), Inches(2.7), Inches(3.6), Inches(3.2),
                     [f"\u25B8  {it}" for it in items], font_size=13, color=LIGHT_GRAY, spacing=Pt(10))

add_accent_bar(slide)
add_slide_number(slide, 19, TOTAL_SLIDES)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 20 — Thank You / Closing
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_shape_rect(slide, Inches(0), Inches(0), Inches(0.15), prs.slide_height, RED_PRIMARY)

add_text_box(slide, Inches(1.5), Inches(1.5), Inches(10), Inches(0.6),
             "KEMENTERIAN KOMUNIKASI DAN DIGITAL", font_size=16, color=SUBTITLE_GRAY, font_name="Segoe UI Light")
add_text_box(slide, Inches(1.5), Inches(2.2), Inches(10), Inches(1.0),
             "Terima Kasih", font_size=54, color=WHITE, bold=True)
add_shape_rect(slide, Inches(1.5), Inches(3.3), Inches(4), Inches(0.04), RED_PRIMARY)
add_text_box(slide, Inches(1.5), Inches(3.6), Inches(10), Inches(0.5),
             "Telkom Insight Hub — Panel Manajemen Data Telekomunikasi", font_size=18, color=GOLD)

# Contact / Sign-off section
add_rounded_rect(slide, Inches(1.5), Inches(4.5), Inches(5), Inches(2.0), DARK_CARD, RED_PRIMARY)
add_text_box(slide, Inches(1.8), Inches(4.65), Inches(4.5), Inches(0.4),
             "Tim Pengembangan", font_size=16, color=RED_PRIMARY, bold=True)
info_lines = [
    "Direktorat Jenderal Penyelenggaraan Pos dan Informatika",
    "Kementerian Komunikasi dan Digital",
    "Republik Indonesia",
    "",
    "etelekomunikasi.komdigi.go.id",
]
add_bullet_frame(slide, Inches(1.8), Inches(5.05), Inches(4.5), Inches(1.4),
                 info_lines, font_size=12, color=LIGHT_GRAY, spacing=Pt(3))

# Sign-off boxes
add_rounded_rect(slide, Inches(7.5), Inches(4.5), Inches(5.3), Inches(2.0), DARK_CARD, RED_PRIMARY)
add_text_box(slide, Inches(7.7), Inches(4.65), Inches(4.9), Inches(0.4),
             "Persetujuan Presentasi", font_size=14, color=RED_PRIMARY, bold=True)
signoff_lines = [
    "Disusun oleh: _________________________",
    "Tanggal: _________________________",
    "",
    "Disetujui oleh: _________________________",
    "Tanggal: _________________________",
]
add_bullet_frame(slide, Inches(7.7), Inches(5.05), Inches(4.9), Inches(1.4),
                 signoff_lines, font_size=11, color=LIGHT_GRAY, spacing=Pt(3))

add_accent_bar(slide)
add_slide_number(slide, 20, TOTAL_SLIDES)

# ── Save ───────────────────────────────────────────────────────────────────
output_path = os.path.join(os.path.dirname(__file__), "PRESENTATION.pptx")
prs.save(output_path)
print(f"Saved: {output_path}")
